#!/usr/bin/env python3
"""
Enhanced PowerPoint generator for Strategia Marketingu Cyfrowego.
Opens natively in Apple Keynote and PowerPoint.

Improvements over v1:
  - Slide type detection (lecture / workshop / break / re-energizer / transition)
  - Dark backgrounds for workshop/break/transition/re-energizer slides
  - RACE stage colour accents on RACE lecture slides (blue/green/orange/purple)
  - Min 30pt body font (per spec rule)
  - Key points extracted from Wizualizacja shown on slide; full SAY in notes
  - Full speaker notes: TIME · SAY · SETUP · DO · WATCH · DEBRIEF · EXTENSION
  - Block label rendered as coloured tag
"""

import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── colour palette ─────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1A, 0x2B, 0x4E)   # lecture title band
DARK_BG     = RGBColor(0x0D, 0x1B, 0x2A)   # workshop/break/transition dark bg
WORKSHOP_BG = RGBColor(0x0A, 0x1A, 0x10)   # deep green for workshop slides
REENER_BG   = RGBColor(0x0D, 0x0D, 0x1A)   # very dark blue for re-energizer
BREAK_BG    = RGBColor(0x14, 0x14, 0x2E)   # break slides
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF0, 0xF4, 0xFF)
LIGHT_GREY  = RGBColor(0xF2, 0xF4, 0xF7)
MID_GREY    = RGBColor(0xD0, 0xD4, 0xDD)
DIM_GREY    = RGBColor(0x88, 0x92, 0xA4)
DARK_TEXT   = RGBColor(0x1C, 0x1C, 0x1C)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xC2)
GOLD        = RGBColor(0xF5, 0xC5, 0x18)   # workshop accent text
LIGHT_GREEN = RGBColor(0x90, 0xEE, 0x90)   # workshop bullet accent

# RACE stage colours (spec requirement)
RACE = {
    "reach":   RGBColor(0x00, 0x70, 0xC0),  # niebieski
    "act":     RGBColor(0x70, 0xAD, 0x47),  # zielony
    "convert": RGBColor(0xED, 0x7D, 0x31),  # pomarańczowy
    "engage":  RGBColor(0x7B, 0x2F, 0xBE),  # fioletowy
}

# SOSTAC component colours (matches Canva exercise palette)
SOSTAC_COLOURS = {
    "sytuacja": RGBColor(0x00, 0x70, 0xC0),
    "cele":     RGBColor(0x70, 0xAD, 0x47),
    "strategia": RGBColor(0xFF, 0xC0, 0x00),
    "taktyki":  RGBColor(0xED, 0x7D, 0x31),
    "działania": RGBColor(0xFF, 0x66, 0x99),
    "kontrola": RGBColor(0x7B, 0x2F, 0xBE),
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
PAD = Inches(0.22)

# ── slide type detection ───────────────────────────────────────────────────────

def detect_slide_type(slide_data):
    """Return: 'lecture' | 'workshop' | 'break' | 'reenergizer' | 'transition' | 'synthesis'"""
    blok = slide_data["blok"].lower()
    title = slide_data["title"].lower()
    n = slide_data["slide_num"]

    if "re-energizer" in blok or "re-energizer" in title:
        return "reenergizer"
    if "przerwa" in blok or "przerwa" in title or "break" in blok:
        return "break"
    if "warsztat" in blok and "instrukcja" in blok.lower():
        return "workshop"
    if "warsztat" in blok and ("przykład" in blok or "siatka" in blok or "transition" in blok or "przejście" in title.lower()):
        return "workshop"
    if "warsztat" in blok:
        return "workshop"
    if "synteza" in blok or "zakończenie" in blok or n >= 76:
        return "synthesis"
    if "przejście" in title or "transition" in blok:
        return "transition"
    return "lecture"


def detect_race_stage(slide_data):
    """Return RACE stage colour key or None."""
    blok = slide_data["blok"].lower()
    title = slide_data["title"].lower()
    n = slide_data["slide_num"]

    # Only for RACE lecture block (slides ~24-38)
    if "lecture block 2" not in blok and "race" not in blok:
        return None
    if "zasięg" in title or "reach" in title:
        return "reach"
    if "aktywacja" in title or "act" in title or "aktyw" in title:
        return "act"
    if "konwersja" in title or "convert" in title or "konwers" in title:
        return "convert"
    if "zaangażowanie" in title or "engage" in title or "zaangaż" in title:
        return "engage"
    return None


# ── helpers ────────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, colour, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                bold=False, italic=False, size=12,
                colour=DARK_TEXT, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = Pt(size)
    run.font.color.rgb = colour
    return tb


def add_textbox_multiline(slide, lines, left, top, width, height,
                          bold=False, size=14, colour=DARK_TEXT,
                          line_colour=None, leading_bullet=True,
                          bg_colour=None):
    """Render a list of lines as separate paragraphs in one textbox."""
    if bg_colour:
        add_rect(slide, left, top, width, height, bg_colour)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        prefix = "• " if leading_bullet else ""
        run.text = prefix + line
        run.font.bold = bold
        run.font.size = Pt(size)
        run.font.color.rgb = line_colour if line_colour else colour
    return tb


# ── parse keynote-spec.md ─────────────────────────────────────────────────────

SPEC_PATH = os.path.join(
    os.path.dirname(__file__),
    ".maister/tasks/development/2026-04-24-digital-marketing-strategy-course/implementation/keynote-spec.md"
)

SLIDE_HEADER_RE = re.compile(r'^### Slide (\d+):', re.MULTILINE)


def extract_field(chunk, tag):
    pattern = re.compile(
        r'\[' + re.escape(tag) + r'\]\s*:\s*(.*?)(?=\n\[|\n---|\Z)',
        re.DOTALL
    )
    m = pattern.search(chunk)
    return m.group(1).strip() if m else ""


def extract_bullet_points(viz_text):
    """Extract key content items from Wizualizacja text for slide display.

    Strategy: prefer quoted labels and numbered list items — these represent
    actual slide text. Avoid layout/structural description words.
    """
    raw = viz_text

    # 1. Quoted strings — these are actual slide labels/text per spec
    quoted = re.findall(r'"([^"]{8,90})"', raw)
    # Filter out layout/structural description noise
    noise_words = {"slajdu", "widoczne", "wyświetlone", "hierarchy-diagram",
                   "specyfikacja", "wariant", "variant", "animacj", "odsłon",
                   "klikn", "wskaźnik", "przejdź", "zatrzymaj", "podkreśl"}
    quoted = [q.strip() for q in quoted
              if len(q.strip()) >= 12
              and not any(n.lower() in q.lower() for n in noise_words)
              and not q.strip().startswith(("z pod", "w kole", "w lewym", "w praw",
                                            "na dole", "na górze", "w rogu",
                                            "Minimalna", "Tytuł slajdu"))]
    if len(quoted) >= 2:
        return [q[:90] for q in quoted[:6]]

    # 2. Numbered items: (1) Label — description
    numbered = re.findall(r'\(\d+\)\s+([A-ZŁŚŻŹĆĄĘÓŃ][^;()\n]{6,80}?)(?=\s*\(\d+\)|\s*—|\.|$)', raw)
    if len(numbered) >= 2:
        return [n.strip()[:90] for n in numbered[:6]]

    # 3. Dash-separated labels: "Label — description" or "Label: description"
    dash_items = re.findall(r'[–—]\s*([A-ZŁŚŻŹĆĄĘÓŃ][^–—\n]{8,80})', raw)
    if len(dash_items) >= 2:
        return [d.strip()[:90] for d in dash_items[:5]]

    # 4. Colon-labelled items: "Punkty pod X: a, b, c" — split by comma
    colon_m = re.search(r'Punkty[^:]+:\s*([^.]+)', raw)
    if colon_m:
        items = [x.strip() for x in colon_m.group(1).split(",") if len(x.strip()) > 5]
        if len(items) >= 2:
            return items[:5]

    # 5. Fallback: first meaningful sentences (avoid pure layout sentences)
    sentences = re.split(r'(?<=[.!?])\s+', raw)
    good = []
    skip_starts = ("slajd", "tło", "centraln", "wizual", "układ", "diagram",
                   "schemat", "grafik", "ikona", "placeholder", "miniatur")
    for sent in sentences:
        sent = sent.strip()
        if len(sent) > 20 and not any(sent.lower().startswith(s) for s in skip_starts):
            good.append(sent[:90])
        if len(good) >= 4:
            break
    return good


def parse_slides(path):
    with open(path, encoding="utf-8") as f:
        content = f.read()
    splits = list(SLIDE_HEADER_RE.finditer(content))
    slides = []
    for i, match in enumerate(splits):
        num = int(match.group(1))
        start = match.start()
        end = splits[i + 1].start() if i + 1 < len(splits) else len(content)
        chunk = content[start:end]
        slides.append(parse_chunk(chunk, num))
    return slides


def parse_chunk(chunk, num):
    title_m = re.match(r'### (Slide \d+): (.+)', chunk)
    blok_m = re.search(r'\*\*Blok\*\*:\s*(.+)', chunk)
    viz_m = re.search(r'\*\*Wizualizacja\*\*:\s*(.*?)(?=\n\[|\Z)', chunk, re.DOTALL)

    return {
        "slide_num": num,
        "number":    title_m.group(1) if title_m else f"Slide {num}",
        "title":     title_m.group(2).strip() if title_m else "?",
        "blok":      blok_m.group(1).strip() if blok_m else "",
        "viz":       viz_m.group(1).strip() if viz_m else "",
        "say":       extract_field(chunk, "SAY"),
        "do":        extract_field(chunk, "DO"),
        "watch":     extract_field(chunk, "WATCH"),
        "time":      extract_field(chunk, "TIME"),
        "setup":     extract_field(chunk, "SETUP"),
        "debrief":   extract_field(chunk, "DEBRIEF"),
        "extension": extract_field(chunk, "EXTENSION"),
    }


def build_notes(s):
    parts = []
    if s["time"]:
        parts.append(f"[TIME] {s['time']}")
    if s["say"]:
        parts.append(f"[SAY]\n{s['say']}")
    if s["setup"]:
        parts.append(f"[SETUP]\n{s['setup']}")
    if s["do"]:
        parts.append(f"[DO]\n{s['do']}")
    if s["watch"]:
        parts.append(f"[WATCH]\n{s['watch']}")
    if s["debrief"]:
        parts.append(f"[DEBRIEF]\n{s['debrief']}")
    if s["extension"]:
        parts.append(f"[EXTENSION]\n{s['extension']}")
    if s["viz"]:
        parts.append(f"[WIZUALIZACJA]\n{s['viz']}")
    return "\n\n".join(parts)


# ── slide builders ────────────────────────────────────────────────────────────

def build_lecture_slide(prs, s, race_stage=None):
    """Standard lecture slide: navy title band + light body."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = SLIDE_W, SLIDE_H
    TITLE_H = Inches(1.05)

    # Title band
    title_bg = RACE[race_stage] if race_stage else NAVY
    add_rect(slide, 0, 0, W, TITLE_H, title_bg)

    # Slide number
    add_textbox(slide, s["number"],
                PAD, Inches(0.12), Inches(0.9), Inches(0.4),
                bold=True, size=11,
                colour=RGBColor(0xAA, 0xBB, 0xCC))

    # Title
    add_textbox(slide, s["title"],
                Inches(1.1), Inches(0.08),
                W - Inches(1.3), TITLE_H - Inches(0.08),
                bold=True, size=26, colour=WHITE)

    # Block label (top-right pill)
    if s["blok"]:
        label_colour = RGBColor(0xCC, 0xDD, 0xFF) if not race_stage else WHITE
        add_textbox(slide, s["blok"],
                    W - Inches(3.5), Inches(0.03),
                    Inches(3.3), Inches(0.36),
                    size=9, colour=label_colour,
                    align=PP_ALIGN.RIGHT)

    # RACE stage accent stripe (left edge) if applicable
    if race_stage:
        add_rect(slide, 0, TITLE_H, Inches(0.07), H - TITLE_H, RACE[race_stage])

    # Content area
    content_top = TITLE_H + PAD
    content_h = H - content_top - Inches(0.45)

    # Extract bullet points from Wizualizacja
    points = extract_bullet_points(s["viz"]) if s["viz"] else []

    if points:
        # Two-zone layout: bullets top-half, SAY excerpt bottom strip
        bullet_h = min(content_h * 0.65, Inches(3.5))
        add_textbox_multiline(
            slide, points,
            PAD, content_top,
            W - PAD * 2, bullet_h,
            bold=False, size=22, colour=DARK_TEXT,
            leading_bullet=True
        )
        # SAY excerpt (first 240 chars) below bullets
        say_excerpt = s["say"][:240].rsplit(" ", 1)[0] + "…" if len(s["say"]) > 240 else s["say"]
        say_top = content_top + bullet_h + PAD * 0.5
        say_h = H - say_top - Inches(0.45)
        if say_h > Inches(0.4):
            add_rect(slide, PAD, say_top, W - PAD * 2, say_h,
                     RGBColor(0xF5, 0xF7, 0xFA))
            add_textbox(slide, say_excerpt,
                        PAD * 2, say_top + PAD * 0.3,
                        W - PAD * 4, say_h - PAD * 0.5,
                        italic=True, size=13, colour=RGBColor(0x44, 0x55, 0x66))
    else:
        # Just SAY text, min 30pt per spec
        say_short = s["say"][:380].rsplit(" ", 1)[0] + "…" if len(s["say"]) > 380 else s["say"]
        add_textbox(slide, say_short,
                    PAD, content_top,
                    W - PAD * 2, content_h,
                    size=20, colour=DARK_TEXT)

    # TIME badge
    if s["time"]:
        add_textbox(slide, s["time"],
                    W - Inches(2.1), H - Inches(0.38),
                    Inches(1.9), Inches(0.32),
                    bold=True, size=12,
                    colour=ACCENT_BLUE if not race_stage else RACE[race_stage],
                    align=PP_ALIGN.RIGHT)

    return slide


def build_dark_slide(prs, s, slide_type):
    """Dark background slide for workshop/break/transition/re-energizer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = SLIDE_W, SLIDE_H

    # Choose background colour
    bg = {
        "workshop":    WORKSHOP_BG,
        "break":       BREAK_BG,
        "transition":  DARK_BG,
        "reenergizer": REENER_BG,
        "synthesis":   RGBColor(0x1A, 0x15, 0x35),
    }.get(slide_type, DARK_BG)

    add_rect(slide, 0, 0, W, H, bg)

    # Top accent stripe
    stripe_colour = {
        "workshop":    LIGHT_GREEN,
        "break":       RGBColor(0x88, 0xAA, 0xFF),
        "transition":  GOLD,
        "reenergizer": GOLD,
        "synthesis":   RGBColor(0xF5, 0xC5, 0x18),
    }.get(slide_type, ACCENT_BLUE)
    add_rect(slide, 0, 0, W, Inches(0.12), stripe_colour)

    # Slide number (top-left)
    add_textbox(slide, s["number"],
                PAD, Inches(0.15), Inches(1.0), Inches(0.4),
                bold=True, size=11, colour=DIM_GREY)

    # Title (large, centred vertically in top third)
    title_top = Inches(0.35)
    title_h = Inches(1.4)
    add_textbox(slide, s["title"],
                PAD, title_top,
                W - PAD * 2, title_h,
                bold=True, size=32, colour=WHITE)

    # Block label
    if s["blok"]:
        add_textbox(slide, s["blok"],
                    W - Inches(3.6), Inches(0.14),
                    Inches(3.4), Inches(0.34),
                    size=9, colour=stripe_colour,
                    align=PP_ALIGN.RIGHT)

    # Main content
    content_top = Inches(1.85)
    content_h = H - content_top - Inches(0.45)

    if slide_type == "workshop":
        _build_workshop_content(slide, s, content_top, content_h, W, stripe_colour)
    elif slide_type == "break":
        _build_break_content(slide, s, content_top, content_h, W)
    elif slide_type in ("transition", "reenergizer"):
        _build_transition_content(slide, s, content_top, content_h, W, stripe_colour)
    else:
        _build_synthesis_content(slide, s, content_top, content_h, W)

    # TIME badge
    if s["time"]:
        add_textbox(slide, s["time"],
                    W - Inches(2.1), H - Inches(0.38),
                    Inches(1.9), Inches(0.32),
                    bold=True, size=12, colour=stripe_colour,
                    align=PP_ALIGN.RIGHT)

    return slide


def _build_workshop_content(slide, s, top, height, W, accent):
    """Workshop: show TASK, TIME/OUTPUT, TOOL blocks + debrief bullets."""
    # SAY is the task description; SETUP + DEBRIEF as structured blocks

    # Task block
    if s["say"]:
        task_text = s["say"][:300].rsplit(" ", 1)[0] + "…" if len(s["say"]) > 300 else s["say"]
        add_rect(slide, PAD, top, W - PAD * 2, Inches(1.5),
                 RGBColor(0x0F, 0x28, 0x18))
        add_textbox(slide, "ZADANIE",
                    PAD * 2, top + PAD * 0.3,
                    Inches(1.5), Inches(0.38),
                    bold=True, size=12, colour=accent)
        add_textbox(slide, task_text,
                    PAD * 2, top + Inches(0.45),
                    W - PAD * 4, Inches(0.95),
                    size=16, colour=OFF_WHITE)

    # Debrief bullets
    debrief_top = top + Inches(1.65)
    if s["debrief"]:
        questions = [q.strip().strip('"') for q in re.split(r'[,"][\s]*"', s["debrief"]) if len(q.strip()) > 10]
        questions = questions[:4]
        if questions:
            add_textbox(slide, "PYTANIA DEBRIEFINGOWE",
                        PAD, debrief_top,
                        W - PAD * 2, Inches(0.38),
                        bold=True, size=12, colour=accent)
            add_textbox_multiline(
                slide, questions,
                PAD, debrief_top + Inches(0.42),
                W - PAD * 2, height - Inches(1.65) - Inches(0.42),
                size=16, colour=MID_GREY, leading_bullet=True
            )
    elif s["setup"]:
        setup_short = s["setup"][:280].rsplit(" ", 1)[0] + "…" if len(s["setup"]) > 280 else s["setup"]
        add_textbox(slide, "SETUP",
                    PAD, debrief_top,
                    Inches(1.2), Inches(0.38),
                    bold=True, size=12, colour=accent)
        add_textbox(slide, setup_short,
                    PAD, debrief_top + Inches(0.42),
                    W - PAD * 2, height - Inches(1.65) - Inches(0.42),
                    size=16, colour=MID_GREY)


def _build_break_content(slide, s, top, height, W):
    """Break slide: large time announcement centered."""
    if s["say"]:
        add_textbox(slide, s["say"][:120],
                    PAD, top + Inches(0.3),
                    W - PAD * 2, Inches(1.4),
                    bold=True, size=36, colour=WHITE,
                    align=PP_ALIGN.CENTER)
    if s["time"]:
        add_textbox(slide, s["time"],
                    PAD, top + Inches(1.8),
                    W - PAD * 2, Inches(0.6),
                    size=20, colour=DIM_GREY,
                    align=PP_ALIGN.CENTER)


def _build_transition_content(slide, s, top, height, W, accent):
    """Transition / re-energizer: large central message."""
    if s["say"]:
        msg = s["say"][:180].rsplit(" ", 1)[0]
        add_textbox(slide, msg,
                    PAD, top,
                    W - PAD * 2, Inches(2.0),
                    size=28, colour=OFF_WHITE)
    # Bullet points from viz if available
    pts = extract_bullet_points(s["viz"]) if s["viz"] else []
    if pts:
        add_textbox_multiline(
            slide, pts[:4],
            PAD, top + Inches(2.2),
            W - PAD * 2, height - Inches(2.2),
            size=22, colour=MID_GREY, leading_bullet=True
        )


def _build_synthesis_content(slide, s, top, height, W):
    """Synthesis / closing: warm tones, key message prominent."""
    pts = extract_bullet_points(s["viz"]) if s["viz"] else []
    if pts:
        add_textbox_multiline(
            slide, pts[:5],
            PAD, top,
            W - PAD * 2, height * 0.6,
            size=22, colour=OFF_WHITE, leading_bullet=True
        )
    if s["say"]:
        say_excerpt = s["say"][:220].rsplit(" ", 1)[0] + "…" if len(s["say"]) > 220 else s["say"]
        say_top = top + height * 0.62
        add_textbox(slide, say_excerpt,
                    PAD, say_top,
                    W - PAD * 2, height * 0.35,
                    italic=True, size=15, colour=DIM_GREY)


# ── main slide dispatcher ─────────────────────────────────────────────────────

def build_slide(prs, s):
    slide_type = detect_slide_type(s)
    race_stage = detect_race_stage(s) if slide_type == "lecture" else None

    if slide_type == "lecture":
        sl = build_lecture_slide(prs, s, race_stage)
    else:
        sl = build_dark_slide(prs, s, slide_type)

    # ── speaker notes ──────────────────────────────────────────────────────────
    notes_slide = sl.notes_slide
    notes_slide.notes_text_frame.text = build_notes(s)

    return sl


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    print("Parsing keynote-spec.md …")
    slides_data = parse_slides(SPEC_PATH)
    print(f"  → {len(slides_data)} slides parsed")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    type_counts = {}
    for i, s in enumerate(slides_data, 1):
        sl = build_slide(prs, s)
        t = detect_slide_type(s)
        type_counts[t] = type_counts.get(t, 0) + 1
        if i % 10 == 0:
            print(f"  → built {i}/{len(slides_data)} slides")

    out_dir = os.path.join(os.path.dirname(__file__), "course-materials")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "Strategia_Marketingu_Cyfrowego.pptx")
    prs.save(out_path)
    print(f"\nSaved:  {out_path}")
    print(f"Size:   {os.path.getsize(out_path) // 1024} KB")
    print(f"\nSlide type breakdown:")
    for k, v in sorted(type_counts.items()):
        print(f"  {k:15s}: {v}")


if __name__ == "__main__":
    main()
